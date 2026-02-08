#!/usr/bin/env python3
"""Generate Edwards Project Operation Board Report PPT"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
SCREENSHOTS_DIR = os.path.join(SCRIPT_DIR, "screenshots")
OUTPUT_PATH = os.path.join(SCRIPT_DIR, "Edwards_Project_Operation_Board_Report.pptx")

# Color scheme
EDWARDS_BLUE = RGBColor(0x00, 0x5B, 0xA2)  # Edwards brand blue
DARK_BLUE = RGBColor(0x1A, 0x23, 0x5B)
MEDIUM_BLUE = RGBColor(0x2D, 0x5E, 0x9E)
LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
ACCENT_ORANGE = RGBColor(0xF3, 0x9C, 0x12)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
TABLE_HEADER_BG = RGBColor(0x1A, 0x23, 0x5B)
TABLE_ALT_BG = RGBColor(0xF0, 0xF4, 0xFA)

# Slide dimensions (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    """Set slide background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """Add a rectangle shape"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.fill.solid()
    if line_width:
        shape.line.width = line_width
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=None):
    """Add a rounded rectangle"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    return shape


def set_text(shape, text, font_size=12, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name="Malgun Gothic"):
    """Set text in a shape"""
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(text_frame, text, font_size=12, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, space_before=0, space_after=0, font_name="Malgun Gothic"):
    """Add a paragraph to a text frame"""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    return p


def add_table(slide, rows, cols, left, top, width, height):
    """Add a table to the slide"""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    return table_shape.table


def style_table_cell(cell, text, font_size=10, bold=False, color=BLACK, bg_color=None, alignment=PP_ALIGN.LEFT, font_name="Malgun Gothic"):
    """Style a table cell"""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color
    # Margins
    cell.margin_left = Emu(91440)  # ~0.1 inch
    cell.margin_right = Emu(91440)
    cell.margin_top = Emu(45720)
    cell.margin_bottom = Emu(45720)


def style_header_row(table, headers, font_size=10):
    """Style the header row of a table"""
    for i, header_text in enumerate(headers):
        style_table_cell(
            table.cell(0, i), header_text,
            font_size=font_size, bold=True, color=WHITE,
            bg_color=TABLE_HEADER_BG, alignment=PP_ALIGN.CENTER
        )


def style_data_row(table, row_idx, data, font_size=9, first_col_bold=False):
    """Style a data row"""
    bg = TABLE_ALT_BG if row_idx % 2 == 0 else WHITE
    for i, cell_text in enumerate(data):
        bold = first_col_bold and i == 0
        style_table_cell(
            table.cell(row_idx, i), cell_text,
            font_size=font_size, bold=bold, color=DARK_GRAY, bg_color=bg
        )


def add_section_title(slide, text, left, top, width=Inches(12), font_size=22):
    """Add a section title"""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    tf = set_text(txBox, text, font_size=font_size, bold=True, color=DARK_BLUE)
    # Add underline via shape
    line = add_shape(slide, left, top + Inches(0.5), width, Pt(3), fill_color=EDWARDS_BLUE)
    return txBox


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=14):
    """Add a subtitle"""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    set_text(txBox, text, font_size=font_size, bold=True, color=MEDIUM_BLUE)
    return txBox


def add_body_text(slide, text, left, top, width=Inches(12), height=Inches(0.4), font_size=11, color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT):
    """Add body text"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    set_text(txBox, text, font_size=font_size, color=color, bold=bold, alignment=alignment)
    return txBox


def add_bullet_list(slide, items, left, top, width=Inches(12), height=Inches(2), font_size=11):
    """Add a bullet list"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Malgun Gothic"
        p.space_before = Pt(4)
        p.level = 0
    return txBox


def add_callout_box(slide, text, left, top, width, height, bg_color=LIGHT_BLUE, text_color=DARK_BLUE, font_size=11, bold=False):
    """Add a callout/highlight box"""
    shape = add_rounded_rect(slide, left, top, width, height, fill_color=bg_color)
    set_text(shape, text, font_size=font_size, bold=bold, color=text_color, alignment=PP_ALIGN.LEFT)
    shape.text_frame.margin_left = Emu(182880)
    shape.text_frame.margin_top = Emu(91440)
    return shape


def add_icon_card(slide, title, description, left, top, width, height, icon_text="", bg_color=WHITE, accent_color=EDWARDS_BLUE):
    """Add an icon card with title and description"""
    # Card background
    card = add_rounded_rect(slide, left, top, width, height, fill_color=bg_color)
    card.shadow.inherit = False

    # Accent top bar
    bar = add_shape(slide, left + Emu(45720), top + Emu(45720), width - Emu(91440), Pt(4), fill_color=accent_color)

    # Title
    title_box = slide.shapes.add_textbox(left + Emu(91440), top + Emu(136000), width - Emu(182880), Inches(0.35))
    set_text(title_box, title, font_size=12, bold=True, color=DARK_BLUE)

    # Description
    desc_box = slide.shapes.add_textbox(left + Emu(91440), top + Emu(400000), width - Emu(182880), height - Emu(500000))
    tf = set_text(desc_box, description, font_size=9, color=MEDIUM_GRAY)
    tf.word_wrap = True

    return card


# ============================================================
# SLIDE BUILDERS
# ============================================================

def build_slide_1_cover(prs):
    """Page 1: Cover slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, DARK_BLUE)

    # Top decorative bar
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), fill_color=EDWARDS_BLUE)

    # Left accent bar
    add_shape(slide, Inches(0.8), Inches(1.8), Pt(5), Inches(3.5), fill_color=EDWARDS_BLUE)

    # Main title
    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10), Inches(1.0))
    tf = set_text(txBox, "Edwards Project Operation Board", font_size=38, bold=True, color=WHITE)

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1.2), Inches(2.9), Inches(10), Inches(0.6))
    set_text(txBox2, "Engineering Resource & Project Intelligence Platform", font_size=20, color=RGBColor(0x8E, 0xB8, 0xE5))

    # Description
    txBox3 = slide.shapes.add_textbox(Inches(1.2), Inches(3.8), Inches(10), Inches(0.5))
    set_text(txBox3, "우리 엔지니어링 조직의 리소스를 한 눈에, 한 곳에서", font_size=16, color=RGBColor(0xBB, 0xD5, 0xED))

    # Live badge
    badge = add_rounded_rect(slide, Inches(1.2), Inches(4.7), Inches(3.5), Inches(0.5), fill_color=ACCENT_GREEN)
    set_text(badge, "● Live — 4개 부서 운영 중", font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Footer info
    txBox4 = slide.shapes.add_textbox(Inches(1.2), Inches(6.0), Inches(5), Inches(0.8))
    tf4 = set_text(txBox4, "2026.02", font_size=12, color=RGBColor(0x8E, 0xB8, 0xE5))
    add_paragraph(tf4, "EUV Program IS", font_size=12, color=RGBColor(0x8E, 0xB8, 0xE5))

    # Bottom bar
    add_shape(slide, Inches(0), Inches(7.3), SLIDE_WIDTH, Inches(0.2), fill_color=EDWARDS_BLUE)


def build_slide_2_background(prs):
    """Page 2: Background & Design Principles"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Header bar
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), fill_color=EDWARDS_BLUE)

    # Page title
    add_section_title(slide, "배경 및 설계 원칙", Inches(0.6), Inches(0.3), width=Inches(11))

    # --- Problem Table ---
    add_subtitle(slide, "현재의 문제", Inches(0.6), Inches(1.0))

    headers = ["업무 상황", "현재 방식", "소요 시간"]
    data = [
        ['"이 프로젝트에 SW 엔지니어 몇 명 투입 중이야?"', "Excel 파일 찾기, 버전 확인, 수작업 취합", "반나절"],
        ['"다음 분기에 우리 팀 여유 인력 있어?"', "각 PM한테 개별 문의, 답변 대기", "수일"],
        ['"지난달 NPI에 실제 얼마나 투입됐어?"', "SharePoint 리스트 다운로드, 피벗 테이블 작성", "수시간"],
        ['"리소스 부족하다는 근거 데이터 보여줘"', "근거 자료 없음", "불가능"],
    ]

    tbl = add_table(slide, 5, 3, Inches(0.6), Inches(1.4), Inches(12), Inches(1.9))
    # Set column widths
    tbl.columns[0].width = Inches(5.0)
    tbl.columns[1].width = Inches(4.5)
    tbl.columns[2].width = Inches(2.5)
    style_header_row(tbl, headers)
    for i, row_data in enumerate(data):
        bg = TABLE_ALT_BG if i % 2 == 0 else WHITE
        style_table_cell(tbl.cell(i + 1, 0), row_data[0], font_size=9, color=DARK_GRAY, bg_color=bg)
        style_table_cell(tbl.cell(i + 1, 1), row_data[1], font_size=9, color=DARK_GRAY, bg_color=bg)
        # Highlight the time column
        time_color = ACCENT_RED if row_data[2] == "불가능" else ACCENT_ORANGE
        style_table_cell(tbl.cell(i + 1, 2), row_data[2], font_size=10, bold=True, color=time_color, bg_color=bg, alignment=PP_ALIGN.CENTER)

    # --- Design Principle ---
    add_subtitle(slide, '설계 원칙: "넣는 건 쉽게, 분류는 상세히"', Inches(0.6), Inches(3.5))

    headers2 = ["", "기존 SharePoint (악순환)", "Operation Board (선순환)"]
    data2 = [
        ["입력", "8-10 클릭, 2-3분 소요 → 귀찮음", "2-3 클릭 또는 자연어, 30초 이내"],
        ["입력 후", "블랙홀 (취합 시점 불명)", "즉시 내 대시보드에 결과 반영"],
        ["데이터 활용", "어쩌다 한 번 매니저가 피벗 테이블", "매일 자동 집계, 실시간 리포트"],
        ["프로젝트 반영", "안 됨", "리소스 계획에 직접 연동"],
        ["팀원 반응", '"이거 왜 넣어요?"', '"내 FTE가 바로 보이니까 넣게 됨"'],
    ]

    tbl2 = add_table(slide, 6, 3, Inches(0.6), Inches(3.9), Inches(12), Inches(2.5))
    tbl2.columns[0].width = Inches(2.0)
    tbl2.columns[1].width = Inches(5.0)
    tbl2.columns[2].width = Inches(5.0)
    style_header_row(tbl2, headers2)
    for i, row_data in enumerate(data2):
        bg = TABLE_ALT_BG if i % 2 == 0 else WHITE
        style_table_cell(tbl2.cell(i + 1, 0), row_data[0], font_size=9, bold=True, color=DARK_BLUE, bg_color=bg)
        style_table_cell(tbl2.cell(i + 1, 1), row_data[1], font_size=9, color=ACCENT_RED, bg_color=bg)
        style_table_cell(tbl2.cell(i + 1, 2), row_data[2], font_size=9, color=ACCENT_GREEN, bg_color=bg)

    # Callout
    add_callout_box(
        slide,
        "프로젝트 복잡도 증가  ·  인력 확보의 데이터화  ·  비용 투명성 요구 → 지금 해결해야 하는 이유",
        Inches(0.6), Inches(6.6), Inches(12), Inches(0.5),
        bg_color=RGBColor(0xFE, 0xF3, 0xE2),
        text_color=RGBColor(0xA0, 0x5A, 0x00),
        font_size=10, bold=True
    )

    # Page number
    add_body_text(slide, "2", Inches(12.5), Inches(7.0), width=Inches(0.5), font_size=9, color=MEDIUM_GRAY, alignment=PP_ALIGN.RIGHT)


def build_slide_3_input_features(prs):
    """Page 3: System Overview & Input Features (with worklog screenshot)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), fill_color=EDWARDS_BLUE)
    add_section_title(slide, "시스템 개요 및 핵심 기능 (입력 측)", Inches(0.6), Inches(0.3), width=Inches(11))

    # --- 3 Core Values Cards ---
    cards = [
        ("실시간 가시성", "누가, 어디에, 얼마나 투입되는지\n즉시 확인\nExcel 취합 반나절 → 즉시", EDWARDS_BLUE),
        ("데이터 기반 의사결정", "과거 실적 + 미래 계획을\n한 화면에서 비교\n근거 없음 → FTE 데이터 기반", ACCENT_GREEN),
        ("제로 수작업", "Excel 취합, 피벗 테이블,\n수동 보고서 불필요\n월 1-2일 → 자동 생성", RGBColor(0x8E, 0x44, 0xAD)),
    ]

    card_width = Inches(3.8)
    card_start_x = Inches(0.6)
    card_gap = Inches(0.15)
    for i, (title, desc, accent) in enumerate(cards):
        x = card_start_x + i * (card_width + card_gap)
        add_icon_card(slide, title, desc, x, Inches(1.0), card_width, Inches(1.3), accent_color=accent, bg_color=LIGHT_GRAY)

    # --- Smart Worklog Section ---
    add_subtitle(slide, "Smart Worklog — 일일 업무 기록", Inches(0.6), Inches(2.5))

    # Left side: Features table
    add_body_text(slide, "입력 방법 A: 클릭 방식 (기본)", Inches(0.6), Inches(2.95), font_size=10, bold=True, color=DARK_BLUE, width=Inches(5.5))

    feat_headers = ["기능", "설명"]
    feat_data = [
        ["스마트 디폴트", "주 프로젝트, 최근 Activity 자동 선택"],
        ["지난주 복사", "반복 업무는 원클릭 복사"],
        ["미팅 분류", "의사결정 / 정보공유 / 피드백 / 문제해결 구분"],
        ["긴급업무/출장", "돌발 업무 및 출장 추적 가능"],
    ]

    ft = add_table(slide, 5, 2, Inches(0.6), Inches(3.25), Inches(5.5), Inches(1.65))
    ft.columns[0].width = Inches(1.8)
    ft.columns[1].width = Inches(3.7)
    style_header_row(ft, feat_headers, font_size=9)
    for i, row in enumerate(feat_data):
        style_data_row(ft, i + 1, row, font_size=8, first_col_bold=True)

    # AI input section
    add_body_text(slide, "입력 방법 B: AI 자연어 입력 (보조)", Inches(0.6), Inches(5.05), font_size=10, bold=True, color=DARK_BLUE, width=Inches(5.5))

    ai_box = add_callout_box(
        slide,
        '"오전에 OQC 프로젝트 미팅하고 오후에 NPI 407056 설계 작업함"\n→ AI가 자동 분석: OQC Infra / MEET / 4h + ACM NPI 407056 / DESIGN / 4h',
        Inches(0.6), Inches(5.4), Inches(5.5), Inches(0.8),
        bg_color=RGBColor(0xE8, 0xF5, 0xE9), text_color=RGBColor(0x1B, 0x5E, 0x20), font_size=8
    )

    # Cost classification section
    add_body_text(slide, "비용 자동 분류 (Context-Aware)", Inches(0.6), Inches(6.35), font_size=10, bold=True, color=DARK_BLUE, width=Inches(5.5))

    cost_box = add_callout_box(
        slide,
        "Direct Product 60-70%  ·  Direct Project 10-15%  ·  Indirect 15-20%  ·  Overhead <10%\n→ 프로젝트+활동+시간 입력만으로 자동 분류 (95%+ 정확도)",
        Inches(0.6), Inches(6.65), Inches(5.5), Inches(0.6),
        bg_color=RGBColor(0xF3, 0xE5, 0xF5), text_color=RGBColor(0x4A, 0x14, 0x8C), font_size=8
    )

    # Right side: Worklog screenshot
    add_subtitle(slide, "워크로그 화면", Inches(6.4), Inches(2.5), width=Inches(6))

    wl_screenshot = os.path.join(SCREENSHOTS_DIR, "worklogs.png")
    if os.path.exists(wl_screenshot):
        slide.shapes.add_picture(wl_screenshot, Inches(6.4), Inches(2.95), Inches(6.5), Inches(4.3))

    add_body_text(slide, "3", Inches(12.5), Inches(7.0), width=Inches(0.5), font_size=9, color=MEDIUM_GRAY, alignment=PP_ALIGN.RIGHT)


def build_slide_4_management_features(prs):
    """Page 4: Management & Analysis Features (with dashboard & resource-matrix screenshots)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), fill_color=EDWARDS_BLUE)
    add_section_title(slide, "핵심 기능 (관리/분석 측)", Inches(0.6), Inches(0.3), width=Inches(11))

    # --- Left: Dashboard ---
    add_subtitle(slide, "개인 대시보드 (My Dashboard)", Inches(0.6), Inches(1.0), width=Inches(6))

    db_screenshot = os.path.join(SCREENSHOTS_DIR, "dashboard.png")
    if os.path.exists(db_screenshot):
        slide.shapes.add_picture(db_screenshot, Inches(0.6), Inches(1.4), Inches(6.0), Inches(3.0))

    # Dashboard feature list
    dash_items = [
        "My FTE: 프로젝트별 투입 비율 (계획 vs 실적 비교)",
        "업무 분류: Product vs. Support 업무 비율",
        "과부하 경고: FTE > 1.0이면 빨간색 경고",
        "주간 워크로그: 이번 주 기록 현황 요약",
        "마일스톤: 다가오는 Gate Review 일정",
    ]

    db_list_box = slide.shapes.add_textbox(Inches(0.6), Inches(4.5), Inches(6.0), Inches(1.5))
    tf = db_list_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(dash_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(8)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Malgun Gothic"
        p.space_before = Pt(2)

    # --- Right: Resource Matrix ---
    add_subtitle(slide, "리소스 매트릭스 (Resource Matrix)", Inches(6.8), Inches(1.0), width=Inches(6))

    rm_screenshot = os.path.join(SCREENSHOTS_DIR, "resource-matrix.png")
    if os.path.exists(rm_screenshot):
        slide.shapes.add_picture(rm_screenshot, Inches(6.8), Inches(1.4), Inches(6.0), Inches(3.0))

    # Resource Matrix description
    rm_items = [
        "행(Y축): 팀원 이름 / 열(X축): 월별 (12개월+)",
        "셀: 프로젝트별 FTE (색상 코딩)",
        "TBD 포지션: 미배정 자리 → 채용 필요 시점 식별",
        "과배정 식별: FTE > 1.0인 팀원 즉시 파악",
    ]

    rm_list_box = slide.shapes.add_textbox(Inches(6.8), Inches(4.5), Inches(6.0), Inches(1.2))
    tf2 = rm_list_box.text_frame
    tf2.word_wrap = True
    for i, item in enumerate(rm_items):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(8)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Malgun Gothic"
        p.space_before = Pt(2)

    # --- Bottom: Forecasting example ---
    add_subtitle(slide, "리소스 포캐스팅 (12개월 전방 계획) — 예시", Inches(0.6), Inches(5.9), width=Inches(12))

    fc_headers = ["이름", "4월", "5월", "6월", "상태"]
    fc_data = [
        ["김철수", "0.8", "0.8", "0.5", "정상"],
        ["이영희", "1.0", "1.0", "1.0", "주의"],
        ["박준형", "0.6", "0.9", "1.2", "경고"],
        ["TBD-SW Eng", "0.5", "0.8", "1.0", "미배정"],
        ["Gap", "-0.1", "-0.5", "-0.7", "인력 부족"],
    ]

    fc_tbl = add_table(slide, 6, 5, Inches(0.6), Inches(6.25), Inches(7.0), Inches(1.15))
    fc_tbl.columns[0].width = Inches(1.6)
    fc_tbl.columns[1].width = Inches(1.1)
    fc_tbl.columns[2].width = Inches(1.1)
    fc_tbl.columns[3].width = Inches(1.1)
    fc_tbl.columns[4].width = Inches(2.1)
    style_header_row(fc_tbl, fc_headers, font_size=9)
    for i, row in enumerate(fc_data):
        bg = TABLE_ALT_BG if i % 2 == 0 else WHITE
        for j, cell_text in enumerate(row):
            txt_color = DARK_GRAY
            b = False
            if j == 4:  # status column
                if cell_text == "경고":
                    txt_color = ACCENT_RED
                    b = True
                elif cell_text == "주의":
                    txt_color = ACCENT_ORANGE
                    b = True
                elif cell_text == "인력 부족":
                    txt_color = ACCENT_RED
                    b = True
            if i == 4:  # Gap row
                b = True
                if cell_text.startswith("-"):
                    txt_color = ACCENT_RED
            style_table_cell(fc_tbl.cell(i + 1, j), cell_text, font_size=8, bold=b, color=txt_color, bg_color=bg, alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

    # Callout for gap
    add_callout_box(
        slide,
        "→ 6월 Gap -0.7 = 지금 채용을 시작하지 않으면 3개월 후 프로젝트 일정 지연",
        Inches(7.8), Inches(6.25), Inches(5.0), Inches(0.4),
        bg_color=RGBColor(0xFD, 0xED, 0xED), text_color=ACCENT_RED, font_size=9, bold=True
    )

    add_body_text(slide, "4", Inches(12.5), Inches(7.0), width=Inches(0.5), font_size=9, color=MEDIUM_GRAY, alignment=PP_ALIGN.RIGHT)


def build_slide_5_pm_value(prs):
    """Page 5: Core value for PMs"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), fill_color=EDWARDS_BLUE)
    add_section_title(slide, "PM을 위한 핵심 가치 — 데이터 기반 리소스 증명", Inches(0.6), Inches(0.3), width=Inches(11))

    # --- Before/After Comparison ---
    add_subtitle(slide, "기존 방식 vs Operation Board", Inches(0.6), Inches(1.0))

    cmp_headers = ["", "기존", "Operation Board"]
    cmp_data = [
        ["리소스 요청", '"팀이 바빠요. 사람이 부족해요."', '"지난 6개월 평균 FTE 1.15, 월 팀원당 12h 초과 근무"'],
        ["매니저 반응", '"다 바쁘다고 하는데, 근거가 뭐야?"', '"이 데이터면 충분하네. 채용 시작하자."'],
        ["데이터 준비", "Excel 정리 수일 소요", "화면 공유 즉시"],
    ]

    cmp_tbl = add_table(slide, 4, 3, Inches(0.6), Inches(1.4), Inches(12), Inches(1.5))
    cmp_tbl.columns[0].width = Inches(2.0)
    cmp_tbl.columns[1].width = Inches(5.0)
    cmp_tbl.columns[2].width = Inches(5.0)
    style_header_row(cmp_tbl, cmp_headers)
    for i, row in enumerate(cmp_data):
        bg = TABLE_ALT_BG if i % 2 == 0 else WHITE
        style_table_cell(cmp_tbl.cell(i + 1, 0), row[0], font_size=9, bold=True, color=DARK_BLUE, bg_color=bg)
        style_table_cell(cmp_tbl.cell(i + 1, 1), row[1], font_size=9, color=ACCENT_RED, bg_color=bg)
        style_table_cell(cmp_tbl.cell(i + 1, 2), row[2], font_size=9, color=ACCENT_GREEN, bg_color=bg)

    # --- PM Data Table ---
    add_subtitle(slide, "PM이 즉시 활용할 수 있는 데이터", Inches(0.6), Inches(3.15))

    pm_headers = ["데이터", "용도"]
    pm_data = [
        ["Planned vs Actual FTE 트렌드 (6개월)", "지속적 과부하 증명 → 인력 추가 요청 근거"],
        ["FTE > 1.0 초과 빈도 히트맵", "초과 근무가 특정 시기/프로젝트에 집중되는 패턴 파악"],
        ["미배정 TBD 포지션 목록", "채용 필요 시점과 역할 명확화"],
        ["비용 버킷별 비율", "Direct Product 비율로 팀의 핵심 R&D 기여도 증명"],
        ["월별 Capacity Gap", "미래 인력 부족 시점 사전 예측"],
    ]

    pm_tbl = add_table(slide, 6, 2, Inches(0.6), Inches(3.5), Inches(12), Inches(2.2))
    pm_tbl.columns[0].width = Inches(4.5)
    pm_tbl.columns[1].width = Inches(7.5)
    style_header_row(pm_tbl, pm_headers)
    for i, row in enumerate(pm_data):
        bg = TABLE_ALT_BG if i % 2 == 0 else WHITE
        style_table_cell(pm_tbl.cell(i + 1, 0), row[0], font_size=9, bold=True, color=DARK_BLUE, bg_color=bg)
        style_table_cell(pm_tbl.cell(i + 1, 1), row[1], font_size=9, color=DARK_GRAY, bg_color=bg)

    # Key message callout
    msg_box = add_callout_box(
        slide,
        '"감"이 아니라 "팩트"로 이야기할 수 있게 됩니다.\n데이터 없는 리소스 요청은 승인되지 않지만, 데이터 있는 리소스 요청은 거절하기 어렵습니다.',
        Inches(0.6), Inches(5.95), Inches(12), Inches(0.8),
        bg_color=LIGHT_BLUE, text_color=DARK_BLUE, font_size=13, bold=True
    )

    add_body_text(slide, "5", Inches(12.5), Inches(7.0), width=Inches(0.5), font_size=9, color=MEDIUM_GRAY, alignment=PP_ALIGN.RIGHT)


def build_slide_6_expansion(prs):
    """Page 6: Organization-wide expansion & low adoption barrier"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), fill_color=EDWARDS_BLUE)
    add_section_title(slide, "전체 조직 확산의 필요성 및 도입 부담", Inches(0.6), Inches(0.3), width=Inches(11))

    # --- Network Effect Table ---
    add_subtitle(slide, "왜 전체 조직이 써야 하는가: 네트워크 효과", Inches(0.6), Inches(1.0))

    ne_headers = ["", "현재 (4개 부서)", "전체 조직 도입 시"]
    ne_data = [
        ["리소스 가시성", "우리 팀 안에서만 보임", "Cross-functional 전체 리소스 보임"],
        ["프로젝트 리소스", "일부 프로젝트의 일부 팀만", "모든 프로젝트의 전체 투입 인력"],
        ["FTE 데이터", "부분적 (4개 부서 데이터만)", "조직 전체 Capacity 파악"],
        ["포캐스팅", "우리 팀만 예측 가능", "조직 전체 리소스 수급 예측"],
        ["의사결정 레벨", "팀 레벨", "프로그램/조직 레벨"],
    ]

    ne_tbl = add_table(slide, 6, 3, Inches(0.6), Inches(1.35), Inches(12), Inches(2.2))
    ne_tbl.columns[0].width = Inches(2.5)
    ne_tbl.columns[1].width = Inches(4.75)
    ne_tbl.columns[2].width = Inches(4.75)
    style_header_row(ne_tbl, ne_headers)
    for i, row in enumerate(ne_data):
        bg = TABLE_ALT_BG if i % 2 == 0 else WHITE
        style_table_cell(ne_tbl.cell(i + 1, 0), row[0], font_size=9, bold=True, color=DARK_BLUE, bg_color=bg)
        style_table_cell(ne_tbl.cell(i + 1, 1), row[1], font_size=9, color=MEDIUM_GRAY, bg_color=bg)
        style_table_cell(ne_tbl.cell(i + 1, 2), row[2], font_size=9, bold=True, color=ACCENT_GREEN, bg_color=bg)

    # NPI callout
    add_callout_box(
        slide,
        "NPI 프로젝트에 Control, Mechanical, Electrical 팀이 모두 투입됩니다.\n한 팀만 시스템을 쓰면 1/3만 보입니다. 전체가 쓰면 프로젝트의 전체 그림이 완성됩니다.",
        Inches(0.6), Inches(3.65), Inches(12), Inches(0.65),
        bg_color=RGBColor(0xE8, 0xF5, 0xE9), text_color=RGBColor(0x1B, 0x5E, 0x20), font_size=10, bold=True
    )

    # --- Adoption Barrier ---
    add_subtitle(slide, "도입 부담: 거의 없음", Inches(0.6), Inches(4.5))

    ab_headers = ["항목", "내용"]
    ab_data = [
        ["추가 비용", "없음 (자체 인프라 운영)"],
        ["설치", "없음 (웹 브라우저 URL 접속)"],
        ["팀원 교육", "30분 이내 (워크로그 입력 + 대시보드 확인)"],
        ["팀원 일일 부담", "매일 30초 (워크로그 입력)"],
        ["데이터 이관", "기존 SharePoint 데이터 자동 마이그레이션 지원"],
        ["기존 업무 영향", "SharePoint 병행 가능, 점진적 전환"],
    ]

    ab_tbl = add_table(slide, 7, 2, Inches(0.6), Inches(4.85), Inches(6.0), Inches(2.1))
    ab_tbl.columns[0].width = Inches(2.0)
    ab_tbl.columns[1].width = Inches(4.0)
    style_header_row(ab_tbl, ab_headers, font_size=9)
    for i, row in enumerate(ab_data):
        bg = TABLE_ALT_BG if i % 2 == 0 else WHITE
        style_table_cell(ab_tbl.cell(i + 1, 0), row[0], font_size=9, bold=True, color=DARK_BLUE, bg_color=bg)
        style_table_cell(ab_tbl.cell(i + 1, 1), row[1], font_size=9, color=DARK_GRAY, bg_color=bg)

    # Timeline cards
    add_subtitle(slide, "도입 단계", Inches(7.0), Inches(4.5), width=Inches(5.5))

    timeline = [
        ("Week 1", "계정 생성\n& 팀 구성 설정", EDWARDS_BLUE),
        ("Week 2-3", "워크로그 시작\n(매일 30초)", ACCENT_GREEN),
        ("Week 4+", "대시보드 활용\n& 리소스 계획", RGBColor(0x8E, 0x44, 0xAD)),
    ]

    for i, (week, desc, accent) in enumerate(timeline):
        x = Inches(7.0) + i * Inches(2.0)
        card = add_rounded_rect(slide, x, Inches(4.9), Inches(1.85), Inches(1.5), fill_color=LIGHT_GRAY)
        # Accent top
        add_shape(slide, x + Emu(45720), Inches(4.9) + Emu(45720), Inches(1.75), Pt(4), fill_color=accent)

        week_box = slide.shapes.add_textbox(x + Emu(91440), Inches(5.15), Inches(1.7), Inches(0.3))
        set_text(week_box, week, font_size=11, bold=True, color=accent, alignment=PP_ALIGN.CENTER)

        desc_box = slide.shapes.add_textbox(x + Emu(91440), Inches(5.5), Inches(1.7), Inches(0.7))
        set_text(desc_box, desc, font_size=9, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # Arrow connectors between timeline cards
    for i in range(2):
        x = Inches(8.85) + i * Inches(2.0)
        arr = slide.shapes.add_textbox(x, Inches(5.4), Inches(0.3), Inches(0.3))
        set_text(arr, "→", font_size=16, bold=True, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

    add_body_text(slide, "6", Inches(12.5), Inches(7.0), width=Inches(0.5), font_size=9, color=MEDIUM_GRAY, alignment=PP_ALIGN.RIGHT)


def build_slide_7_capability(prs):
    """Page 7: In-house capability & operational results"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), fill_color=EDWARDS_BLUE)
    add_section_title(slide, "자체 구축 역량 및 운영 성과", Inches(0.6), Inches(0.3), width=Inches(11))

    # --- SaaS vs In-house Comparison ---
    add_subtitle(slide, "외부 SaaS vs 자체 구축 비교", Inches(0.6), Inches(1.0))

    saas_headers = ["", "외부 SaaS (PLM, Jira 등)", "Operation Board"]
    saas_data = [
        ["비용", "라이센스 per user/month\n(50명 기준 연 수천만원)", "자체 인프라\n추가 비용 없음"],
        ["커스터마이징", "제한적, 벤더 의존", "우리 프로세스에 100% 맞춤"],
        ["데이터 소유권", "벤더 클라우드", "자체 서버, 완전한 통제"],
        ["변경 속도", "벤더 일정 의존\n(수주~수개월)", "즉시 반영 가능"],
        ["우리 프로세스", "범용 기능을 억지로 맞춤", "Edwards PCP 프로세스 내장"],
        ["한국 환경", "글로벌 기준", "한국 공휴일, 한글 지원"],
    ]

    saas_tbl = add_table(slide, 7, 3, Inches(0.6), Inches(1.35), Inches(12), Inches(2.6))
    saas_tbl.columns[0].width = Inches(2.5)
    saas_tbl.columns[1].width = Inches(4.75)
    saas_tbl.columns[2].width = Inches(4.75)
    style_header_row(saas_tbl, saas_headers)
    for i, row in enumerate(saas_data):
        bg = TABLE_ALT_BG if i % 2 == 0 else WHITE
        style_table_cell(saas_tbl.cell(i + 1, 0), row[0], font_size=9, bold=True, color=DARK_BLUE, bg_color=bg)
        style_table_cell(saas_tbl.cell(i + 1, 1), row[1], font_size=9, color=MEDIUM_GRAY, bg_color=bg)
        style_table_cell(saas_tbl.cell(i + 1, 2), row[2], font_size=9, bold=True, color=ACCENT_GREEN, bg_color=bg)

    # Self-build callout
    add_callout_box(
        slide,
        "이 시스템은 외주가 아닌 우리 엔지니어링 조직 내부에서 직접 설계하고 구축했습니다.\n우리의 PCP 프로세스, 조직 구조, 프로젝트 유형에 맞춰져 있으며, 이는 우리 조직의 자체 SW 역량을 보여주는 것이기도 합니다.",
        Inches(0.6), Inches(4.05), Inches(12), Inches(0.7),
        bg_color=LIGHT_BLUE, text_color=DARK_BLUE, font_size=10, bold=True
    )

    # --- Operational Results ---
    add_subtitle(slide, "현재 운영 성과", Inches(0.6), Inches(4.95))

    res_headers = ["지표", "기존 (SharePoint/Excel)", "현재 (Operation Board)", "개선"]
    res_data = [
        ["워크로그 입력 시간", "2-3분/건", "30초 이내/건", "80% 단축"],
        ["리소스 현황 파악", "반나절 (Excel 취합)", "즉시 (실시간)", "즉시"],
        ["월말 보고서 작성", "1-2일", "자동 생성", "자동화"],
        ["리소스 포캐스팅", "불가능", "12개월 전방", "신규"],
        ["Cross-project 가시성", "개별 문의", "한 화면", "실시간"],
        ["비용 분류", "수동 또는 미분류", "자동 95%+ 정확도", "자동화"],
    ]

    res_tbl = add_table(slide, 7, 4, Inches(0.6), Inches(5.3), Inches(12), Inches(2.0))
    res_tbl.columns[0].width = Inches(3.0)
    res_tbl.columns[1].width = Inches(3.5)
    res_tbl.columns[2].width = Inches(3.5)
    res_tbl.columns[3].width = Inches(2.0)
    style_header_row(res_tbl, res_headers, font_size=9)
    for i, row in enumerate(res_data):
        bg = TABLE_ALT_BG if i % 2 == 0 else WHITE
        style_table_cell(res_tbl.cell(i + 1, 0), row[0], font_size=8, bold=True, color=DARK_BLUE, bg_color=bg)
        style_table_cell(res_tbl.cell(i + 1, 1), row[1], font_size=8, color=ACCENT_RED, bg_color=bg)
        style_table_cell(res_tbl.cell(i + 1, 2), row[2], font_size=8, color=ACCENT_GREEN, bg_color=bg)
        style_table_cell(res_tbl.cell(i + 1, 3), row[3], font_size=8, bold=True, color=EDWARDS_BLUE, bg_color=bg, alignment=PP_ALIGN.CENTER)

    add_body_text(slide, "7", Inches(12.5), Inches(7.0), width=Inches(0.5), font_size=9, color=MEDIUM_GRAY, alignment=PP_ALIGN.RIGHT)


def build_slide_8_roadmap(prs):
    """Page 8: Roadmap & Proposal"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), fill_color=EDWARDS_BLUE)
    add_section_title(slide, "향후 로드맵 및 제안", Inches(0.6), Inches(0.3), width=Inches(11))

    # --- Roadmap ---
    add_subtitle(slide, "발전 로드맵", Inches(0.6), Inches(1.0))

    road_headers = ["Phase", "시기", "주요 내용"]
    road_data = [
        ["Phase 1 (현재)", "완료", "워크로그 관리, 리소스 매트릭스, 개인 대시보드, 프로젝트/마일스톤 관리, 비용 자동 분류, AI 입력 보조"],
        ["Phase 2 (단기)", "3-6개월", "팀 대시보드 고도화, FTE>1.0 자동 경고, 모바일 최적화, Power BI 연동"],
        ["Phase 3 (중장기)", "6-12개월", "What-if 리소스 시뮬레이션, 투입 패턴 분석, Finance ERP(SAP) 연동"],
        ["Phase 4 (비전)", "12개월+", "AI 기반 리소스 추천, 프로젝트 리스크 자동 감지, 글로벌 사이트 확장"],
    ]

    rd_tbl = add_table(slide, 5, 3, Inches(0.6), Inches(1.35), Inches(12), Inches(1.85))
    rd_tbl.columns[0].width = Inches(2.5)
    rd_tbl.columns[1].width = Inches(1.5)
    rd_tbl.columns[2].width = Inches(8.0)
    style_header_row(rd_tbl, road_headers)
    for i, row in enumerate(road_data):
        bg = TABLE_ALT_BG if i % 2 == 0 else WHITE
        phase_color = ACCENT_GREEN if i == 0 else (EDWARDS_BLUE if i == 1 else (MEDIUM_BLUE if i == 2 else MEDIUM_GRAY))
        style_table_cell(rd_tbl.cell(i + 1, 0), row[0], font_size=9, bold=True, color=phase_color, bg_color=bg)
        style_table_cell(rd_tbl.cell(i + 1, 1), row[1], font_size=9, bold=True, color=DARK_GRAY, bg_color=bg, alignment=PP_ALIGN.CENTER)
        style_table_cell(rd_tbl.cell(i + 1, 2), row[2], font_size=8, color=DARK_GRAY, bg_color=bg)

    # --- Proposal ---
    add_subtitle(slide, "제안: 다음 스텝", Inches(0.6), Inches(3.4))

    prop_headers = ["역할", "요청 사항", "소요"]
    prop_data = [
        ["관심 부서", "2주 파일럿 참여 (계정 생성 → 워크로그 시작 → 효과 체험)", "매일 30초"],
        ["PM", "현재 Excel 리소스 플랜을 시스템에 등록 (1회성) → 이후 실시간 추적", "최초 1회"],
        ["FM", "팀원 FTE 현황 확인, TBD 포지션 관리 시작", "주 1회 확인"],
    ]

    pr_tbl = add_table(slide, 4, 3, Inches(0.6), Inches(3.75), Inches(12), Inches(1.2))
    pr_tbl.columns[0].width = Inches(1.5)
    pr_tbl.columns[1].width = Inches(8.5)
    pr_tbl.columns[2].width = Inches(2.0)
    style_header_row(pr_tbl, prop_headers)
    for i, row in enumerate(prop_data):
        bg = TABLE_ALT_BG if i % 2 == 0 else WHITE
        style_table_cell(pr_tbl.cell(i + 1, 0), row[0], font_size=9, bold=True, color=DARK_BLUE, bg_color=bg)
        style_table_cell(pr_tbl.cell(i + 1, 1), row[1], font_size=9, color=DARK_GRAY, bg_color=bg)
        style_table_cell(pr_tbl.cell(i + 1, 2), row[2], font_size=9, bold=True, color=ACCENT_GREEN, bg_color=bg, alignment=PP_ALIGN.CENTER)

    # Support items
    add_subtitle(slide, "지원 사항", Inches(0.6), Inches(5.15), width=Inches(6))
    support_items = [
        "• 30분 온보딩 세션 (팀별 제공)",
        "• 기존 SharePoint 데이터 자동 마이그레이션",
        "• 지속적 기능 개선 및 피드백 반영",
    ]
    sup_box = slide.shapes.add_textbox(Inches(0.6), Inches(5.5), Inches(6), Inches(1.0))
    tf = sup_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(support_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Malgun Gothic"
        p.space_before = Pt(4)

    # Key message callout
    add_callout_box(
        slide,
        "강제가 아닙니다. 2주만 써보시면 됩니다.\n효과가 없으시면 기존 방식으로 돌아가시면 됩니다.\n다만, 써보신 4개 부서에서는 아직 돌아간 분이 없습니다.",
        Inches(6.8), Inches(5.15), Inches(6.0), Inches(1.0),
        bg_color=RGBColor(0xE8, 0xF5, 0xE9), text_color=RGBColor(0x1B, 0x5E, 0x20), font_size=12, bold=True
    )

    # Footer
    footer_box = add_shape(slide, Inches(0), Inches(6.5), SLIDE_WIDTH, Inches(1.0), fill_color=DARK_BLUE)
    ft_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.65), Inches(12), Inches(0.6))
    tf = set_text(ft_box, "접속 정보:  시스템 URL  |  데모 계정: guest@edwards.com  |  담당자: EUV Program IS", font_size=12, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_body_text(slide, "8", Inches(12.5), Inches(7.0), width=Inches(0.5), font_size=9, color=RGBColor(0x88, 0x99, 0xBB), alignment=PP_ALIGN.RIGHT)


def main():
    prs = Presentation()

    # Set 16:9 slide size
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    print("Building slides...")

    print("  Page 1: Cover")
    build_slide_1_cover(prs)

    print("  Page 2: Background & Design Principles")
    build_slide_2_background(prs)

    print("  Page 3: System Overview & Input Features")
    build_slide_3_input_features(prs)

    print("  Page 4: Management & Analysis Features")
    build_slide_4_management_features(prs)

    print("  Page 5: PM Core Value")
    build_slide_5_pm_value(prs)

    print("  Page 6: Organization Expansion")
    build_slide_6_expansion(prs)

    print("  Page 7: In-house Capability")
    build_slide_7_capability(prs)

    print("  Page 8: Roadmap & Proposal")
    build_slide_8_roadmap(prs)

    prs.save(OUTPUT_PATH)
    print(f"\nPPT saved to: {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
