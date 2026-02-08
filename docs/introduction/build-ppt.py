#!/usr/bin/env python3
"""Build PPT from AC template + HTML slide screenshots"""

import os
from pptx import Presentation
from pptx.util import Inches, Emu

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_PATH = os.path.join(SCRIPT_DIR, "ac.templete.pptx")
SLIDES_DIR = os.path.join(SCRIPT_DIR, "slide-screenshots")
OUTPUT_PATH = os.path.join(SCRIPT_DIR, "Edwards_Project_Operation_Board_Report_v3.pptx")

# Slide dimensions (16:9 widescreen)
SLIDE_W = 12192000  # EMU
SLIDE_H = 6858000   # EMU


def main():
    # Open the template
    prs = Presentation(TEMPLATE_PATH)

    # Find the Blank layout (Layout 14)
    blank_layout = None
    for layout in prs.slide_layouts:
        if layout.name == "Blank":
            blank_layout = layout
            break

    if blank_layout is None:
        print("Warning: 'Blank' layout not found, using first layout")
        blank_layout = prs.slide_layouts[0]
    else:
        print(f"Using layout: '{blank_layout.name}'")

    # Remove existing slides (template has 3 sample slides)
    from pptx.oxml.ns import qn
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn('r:id'))
        if rId:
            prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)

    print(f"Cleared template slides. Remaining: {len(prs.slides)}")

    # Add 8 slides with our screenshots
    for i in range(1, 9):
        slide_img = os.path.join(SLIDES_DIR, f"slide{i}.png")
        if not os.path.exists(slide_img):
            print(f"  Warning: {slide_img} not found, skipping")
            continue

        slide = prs.slides.add_slide(blank_layout)

        # Insert the screenshot as a full-page image
        slide.shapes.add_picture(
            slide_img,
            left=0,
            top=0,
            width=SLIDE_W,
            height=SLIDE_H,
        )
        print(f"  Added slide {i}: {slide_img}")

    prs.save(OUTPUT_PATH)
    print(f"\nPPT saved: {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")
    print(f"File size: {os.path.getsize(OUTPUT_PATH) / 1024 / 1024:.1f} MB")


if __name__ == "__main__":
    main()
