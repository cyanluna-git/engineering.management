#!/usr/bin/env python3
"""Generate Edwards Project Operation Board Report PPT — Editable elements version
Uses AC template color scheme + v3 English content with all elements editable in PowerPoint.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_PATH = os.path.join(SCRIPT_DIR, "ac.templete.pptx")
SCREENSHOTS_DIR = os.path.join(SCRIPT_DIR, "screenshots")
OUTPUT_PATH = os.path.join(SCRIPT_DIR, "Edwards_Project_Operation_Board_Report_v4_editable.pptx")

# Atlas Copco Color Scheme
PRIMARY = RGBColor(0x05, 0x4E, 0x5A)
PRIMARY_L = RGBColor(0x0A, 0x6E, 0x7E)
GOLD = RGBColor(0xE1, 0xB7, 0x7E)
GOLD_DARK = RGBColor(0xC8, 0x9B, 0x5E)
SAGE = RGBColor(0x5D, 0x78, 0x75)
NAVY = RGBColor(0x12, 0x3F, 0x6D)
NAVY_DARK = RGBColor(0x06, 0x31, 0x5B)
CORAL = RGBColor(0xF6, 0x83, 0x63)
SUCCESS = RGBColor(0x27, 0xAE, 0x60)
DANGER = RGBColor(0xE7, 0x4C, 0x3C)
WARNING = RGBColor(0xF3, 0x9C, 0x12)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
TEXT = RGBColor(0x2C, 0x3E, 0x50)
TEXT_LIGHT = RGBColor(0x5A, 0x6A, 0x7A)
GRAY = RGBColor(0xA1, 0xA9, 0xB4)
GRAY_LIGHT = RGBColor(0xDF, 0xE4, 0xE7)
GRAY_BG = RGBColor(0xF4, 0xF6, 0xF8)

FONT = "Segoe UI"
FONT_BOLD = "Segoe UI Semibold"

# Slide dimensions (16:9)
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.fill.solid()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.fill.solid()
    return shape


def set_text(shape, text, size=12, bold=False, color=TEXT, align=PP_ALIGN.LEFT, font=FONT):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return tf


def add_para(tf, text, size=12, bold=False, color=TEXT, align=PP_ALIGN.LEFT, font=FONT, space_before=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    return p


def textbox(slide, text, left, top, width, height=Inches(0.4), size=12, bold=False, color=TEXT, align=PP_ALIGN.LEFT, font=FONT):
    box = slide.shapes.add_textbox(left, top, width, height)
    set_text(box, text, size=size, bold=bold, color=color, align=align, font=font)
    box.text_frame.word_wrap = True
    return box


def add_table(slide, rows, cols, left, top, width, height):
    return slide.shapes.add_table(rows, cols, left, top, width, height).table


def cell(table, r, c, text, size=11, bold=False, color=TEXT, bg=None, align=PP_ALIGN.LEFT):
    cl = table.cell(r, c)
    cl.text = ""
    p = cl.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT
    p.alignment = align
    cl.vertical_anchor = MSO_ANCHOR.MIDDLE
    if bg:
        cl.fill.solid()
        cl.fill.fore_color.rgb = bg
    cl.margin_left = Emu(91440)
    cl.margin_right = Emu(91440)
    cl.margin_top = Emu(45720)
    cl.margin_bottom = Emu(45720)


def header_row(table, headers, size=11):
    for i, h in enumerate(headers):
        cell(table, 0, i, h, size=size, bold=True, color=WHITE, bg=PRIMARY, align=PP_ALIGN.CENTER)


def data_row(table, r, data, size=10, colors=None, bolds=None):
    bg = GRAY_BG if r % 2 == 1 else WHITE
    for i, d in enumerate(data):
        c = colors[i] if colors else TEXT
        b = bolds[i] if bolds else False
        cell(table, r, i, d, size=size, bold=b, color=c, bg=bg)


def top_bar(slide):
    add_shape(slide, Emu(0), Emu(0), SLIDE_W, Emu(54000), fill_color=PRIMARY)


def page_title(slide, text, y=Emu(280000)):
    textbox(slide, text, Emu(500000), y, Emu(10000000), Emu(500000),
            size=24, bold=True, color=NAVY_DARK, font=FONT_BOLD)
    add_shape(slide, Emu(500000), y + Emu(440000), Emu(700000), Emu(36000), fill_color=GOLD)


def subtitle(slide, text, left, top, width=Emu(5000000)):
    return textbox(slide, text, left, top, width, Emu(350000),
                   size=16, bold=True, color=PRIMARY, font=FONT_BOLD)


def page_num(slide, num):
    textbox(slide, str(num), Emu(11500000), Emu(6500000), Emu(400000), Emu(250000),
            size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def footer_line(slide):
    add_shape(slide, Emu(500000), Emu(6400000), Emu(11200000), Emu(9000), fill_color=GRAY_LIGHT)


def callout(slide, text, left, top, width, height, bg=GRAY_BG, color=TEXT, size=11, bold=False, border_color=None):
    shape = add_rounded_rect(slide, left, top, width, height, fill_color=bg, line_color=border_color)
    tf = set_text(shape, text, size=size, bold=bold, color=color)
    tf.word_wrap = True
    shape.text_frame.margin_left = Emu(150000)
    shape.text_frame.margin_top = Emu(80000)
    shape.text_frame.margin_right = Emu(150000)
    return shape


# ============================================================
# SLIDE BUILDERS
# ============================================================

def build_slide_1(prs, blank):
    """Cover"""
    slide = prs.slides.add_slide(blank)

    # Left dark panel
    add_shape(slide, Emu(0), Emu(0), Emu(5800000), SLIDE_H, fill_color=PRIMARY)
    # Gold left accent
    add_shape(slide, Emu(0), Emu(0), Emu(50000), Emu(4000000), fill_color=GOLD)

    # Edwards Vacuum label
    textbox(slide, "EDWARDS VACUUM", Emu(500000), Emu(700000), Emu(4500000), Emu(250000),
            size=11, bold=True, color=GOLD, font=FONT_BOLD)

    # Main title
    title_box = slide.shapes.add_textbox(Emu(500000), Emu(1100000), Emu(4800000), Emu(1500000))
    tf = set_text(title_box, "Project", size=36, bold=True, color=WHITE, font=FONT_BOLD)
    add_para(tf, "Operation", size=36, bold=True, color=WHITE, font=FONT_BOLD)
    add_para(tf, "Board", size=36, bold=True, color=WHITE, font=FONT_BOLD)

    # Gold line
    add_shape(slide, Emu(500000), Emu(2700000), Emu(500000), Emu(27000), fill_color=GOLD)

    # Subtitle
    textbox(slide, "Engineering Resource &\nProject Intelligence Platform",
            Emu(500000), Emu(2850000), Emu(4500000), Emu(600000),
            size=15, color=RGBColor(0xBB, 0xD5, 0xDD))

    # Live badge
    badge = add_rounded_rect(slide, Emu(500000), Emu(3600000), Emu(3200000), Emu(380000),
                             fill_color=RGBColor(0x0E, 0x5C, 0x35))
    set_text(badge, "  Live — 4 Divisions in Production", size=12, bold=True, color=RGBColor(0x4A, 0xDE, 0x80),
             align=PP_ALIGN.LEFT)

    # Stats
    stats = [("50+", "Active Users"), ("30s", "Daily Input"), ("12mo", "Forecast")]
    for i, (val, label) in enumerate(stats):
        x = Emu(500000 + i * 1600000)
        textbox(slide, val, x, Emu(4200000), Emu(1200000), Emu(450000),
                size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER, font=FONT_BOLD)
        textbox(slide, label, x, Emu(4600000), Emu(1200000), Emu(250000),
                size=10, color=RGBColor(0x88, 0xAA, 0xBB), align=PP_ALIGN.CENTER)

    # Date / team
    textbox(slide, "February 2026", Emu(500000), Emu(5800000), Emu(3000000), Emu(200000),
            size=10, color=RGBColor(0x88, 0xAA, 0xBB))
    textbox(slide, "EUV Program IS", Emu(500000), Emu(6000000), Emu(3000000), Emu(200000),
            size=10, color=RGBColor(0x88, 0xAA, 0xBB))

    # Right panel — What We Deliver
    textbox(slide, "WHAT WE DELIVER", Emu(6200000), Emu(400000), Emu(5000000), Emu(300000),
            size=11, bold=True, color=GRAY, font=FONT_BOLD)

    cards = [
        ("Real-time Resource Visibility", "See who is working on what project and FTE allocation — updated in real-time.", PRIMARY),
        ("Data-Driven Decision Making", "Justify resource requests with 6-month FTE trends, capacity gaps, and cost data.", GOLD_DARK),
        ("Zero Manual Aggregation", "Eliminate monthly Excel consolidation — saving 1-2 days per month per manager.", SAGE),
        ("Automated Cost Classification", "Rule engine classifies every worklog into cost buckets at 95%+ accuracy.", CORAL),
        ("12-Month Resource Forecasting", "Plan monthly FTE per project per role. Identify hiring gaps with TBD positions.", NAVY),
    ]
    for i, (title, desc, accent) in enumerate(cards):
        y = Emu(800000 + i * 1150000)
        card = add_rounded_rect(slide, Emu(6200000), y, Emu(5500000), Emu(1000000), fill_color=WHITE)
        add_shape(slide, Emu(6200000), y, Emu(40000), Emu(1000000), fill_color=accent)
        textbox(slide, title, Emu(6450000), y + Emu(100000), Emu(5000000), Emu(300000),
                size=14, bold=True, color=NAVY_DARK, font=FONT_BOLD)
        textbox(slide, desc, Emu(6450000), y + Emu(430000), Emu(5000000), Emu(500000),
                size=11, color=TEXT_LIGHT)


def build_slide_2(prs, blank):
    """Background & Design Principles"""
    slide = prs.slides.add_slide(blank)
    top_bar(slide)
    page_title(slide, "Background & Design Principles")

    L = Emu(500000)
    MID = Emu(6200000)
    W_HALF = Emu(5500000)

    # Left: Current Problem
    subtitle(slide, "The Current Problem", L, Emu(850000))
    t = add_table(slide, 5, 3, L, Emu(1200000), W_HALF, Emu(1900000))
    t.columns[0].width = Emu(2300000)
    t.columns[1].width = Emu(2100000)
    t.columns[2].width = Emu(1100000)
    header_row(t, ["Situation", "Current Method", "Time"])
    rows = [
        ['"How many SW engineers?"', "Find Excel, manual merge", "Half a day"],
        ['"Available capacity next Q?"', "Ask each PM, wait", "Several days"],
        ['"NPI effort last month?"', "SharePoint pivot table", "Hours"],
        ['"Data proving shortage?"', "No data available", "Impossible"],
    ]
    time_colors = [WARNING, WARNING, WARNING, DANGER]
    for i, r in enumerate(rows):
        data_row(t, i + 1, r, size=10,
                 colors=[TEXT, TEXT, time_colors[i]],
                 bolds=[False, False, True])

    # Left: Why Now
    subtitle(slide, "Why This Must Be Solved Now", L, Emu(3300000))
    why_items = [
        ("01", "Increasing Project Complexity", "NPI, ETO, CIP running simultaneously. Engineers on 2-3 projects.", CORAL),
        ("02", "Data-Driven Headcount Justification", "Without FTE evidence, headcount requests are routinely rejected.", GOLD_DARK),
        ("03", "Cost Transparency Requirements", "Direct vs. Indirect costs for GAAP compliance.", PRIMARY),
    ]
    for i, (num, title, desc, accent) in enumerate(why_items):
        y = Emu(3650000 + i * 700000)
        box = add_rounded_rect(slide, L, y, W_HALF, Emu(600000), fill_color=GRAY_BG)
        add_shape(slide, L, y, Emu(30000), Emu(600000), fill_color=accent)
        textbox(slide, num, L + Emu(80000), y + Emu(80000), Emu(350000), Emu(250000),
                size=11, bold=True, color=accent, font=FONT_BOLD)
        textbox(slide, title, L + Emu(400000), y + Emu(60000), Emu(4800000), Emu(250000),
                size=12, bold=True, color=NAVY_DARK, font=FONT_BOLD)
        textbox(slide, desc, L + Emu(400000), y + Emu(310000), Emu(4800000), Emu(280000),
                size=10, color=TEXT_LIGHT)

    # Right: Design Principle
    subtitle(slide, 'Design: "Easy to Input, Rich in Classification"', MID, Emu(850000), width=Emu(5800000))
    t2 = add_table(slide, 6, 3, MID, Emu(1200000), Emu(5600000), Emu(2400000))
    t2.columns[0].width = Emu(900000)
    t2.columns[1].width = Emu(2350000)
    t2.columns[2].width = Emu(2350000)
    cell(t2, 0, 0, "", size=10, bold=True, color=WHITE, bg=PRIMARY, align=PP_ALIGN.CENTER)
    cell(t2, 0, 1, "SharePoint (Vicious)", size=10, bold=True, color=WHITE, bg=RGBColor(0x7A, 0x2E, 0x2E), align=PP_ALIGN.CENTER)
    cell(t2, 0, 2, "Operation Board (Virtuous)", size=10, bold=True, color=WHITE, bg=RGBColor(0x1A, 0x6B, 0x40), align=PP_ALIGN.CENTER)
    cmp = [
        ["Input", "8-10 clicks, 2-3 min", "2-3 clicks, under 30 sec"],
        ["After Input", "Black hole — unknown", "Instantly on dashboard"],
        ["Data Use", "Occasional pivot tables", "Daily auto-aggregation"],
        ["Project Link", "Not connected", "Linked to resource plans"],
        ["Team Response", '"Why do I need to log?"', '"I can see my FTE"'],
    ]
    for i, r in enumerate(cmp):
        bg = GRAY_BG if i % 2 == 0 else WHITE
        cell(t2, i + 1, 0, r[0], size=10, bold=True, color=NAVY_DARK, bg=bg)
        cell(t2, i + 1, 1, r[1], size=10, color=DANGER, bg=bg)
        cell(t2, i + 1, 2, r[2], size=10, color=SUCCESS, bg=bg)

    # Cycle visualization
    for i, (title, items, color, bg_c) in enumerate([
        ("Vicious Cycle", "Hard to input → Low participation → Incomplete data → No output → \"Why bother?\"",
         DANGER, RGBColor(0xFD, 0xED, 0xED)),
        ("Virtuous Cycle", "Easy input (30s) → High participation → Rich data → Valuable dashboards → \"I want this\"",
         SUCCESS, RGBColor(0xE8, 0xF5, 0xE9)),
    ]):
        x = MID + Emu(i * 2900000)
        box = add_rounded_rect(slide, x, Emu(3800000), Emu(2700000), Emu(2300000), fill_color=bg_c)
        textbox(slide, title, x + Emu(100000), Emu(3900000), Emu(2500000), Emu(300000),
                size=13, bold=True, color=color, align=PP_ALIGN.CENTER, font=FONT_BOLD)
        textbox(slide, items, x + Emu(100000), Emu(4250000), Emu(2500000), Emu(1700000),
                size=11, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    footer_line(slide)
    page_num(slide, 2)


def build_slide_3(prs, blank):
    """Input Features"""
    slide = prs.slides.add_slide(blank)
    top_bar(slide)
    page_title(slide, "System Overview — Input Features")

    L = Emu(500000)

    # 3 Value cards
    card_colors = [PRIMARY, NAVY, SAGE]
    card_titles = ["Real-time Visibility", "Data-Driven Decisions", "Zero Manual Work"]
    card_descs = ["Excel (half day) → Instant access", "No evidence → FTE-backed arguments", "1-2 days report → Fully automated"]
    for i in range(3):
        x = L + Emu(i * 3900000)
        card = add_rounded_rect(slide, x, Emu(850000), Emu(3700000), Emu(650000), fill_color=card_colors[i])
        textbox(slide, card_titles[i], x + Emu(150000), Emu(900000), Emu(3400000), Emu(300000),
                size=13, bold=True, color=WHITE, font=FONT_BOLD)
        textbox(slide, card_descs[i], x + Emu(150000), Emu(1170000), Emu(3400000), Emu(250000),
                size=11, color=RGBColor(0xDD, 0xDD, 0xDD))

    # Left: Feature tables
    LEFT_W = Emu(5500000)
    subtitle(slide, "Feature 1: Smart Worklog", L, Emu(1650000))

    textbox(slide, "Method A: Click-based Entry", L, Emu(1980000), Emu(5000000), Emu(250000),
            size=12, bold=True, color=NAVY_DARK, font=FONT_BOLD)

    t = add_table(slide, 5, 2, L, Emu(2250000), LEFT_W, Emu(1600000))
    t.columns[0].width = Emu(1600000)
    t.columns[1].width = Emu(3900000)
    header_row(t, ["Feature", "Description"])
    feat = [
        ["Smart Defaults", "Auto-selects primary project & recent activities"],
        ["Copy Last Week", "One-click copy for recurring tasks"],
        ["Meeting Types", "Decision / Info-sharing / Feedback / Problem-solving"],
        ["Flags", "Urgent task & business trip flags"],
    ]
    for i, r in enumerate(feat):
        data_row(t, i + 1, r, size=10, bolds=[True, False], colors=[NAVY_DARK, TEXT])

    # AI input
    textbox(slide, "Method B: AI Natural Language Input", L, Emu(3950000), Emu(5000000), Emu(250000),
            size=12, bold=True, color=NAVY_DARK, font=FONT_BOLD)
    callout(slide,
            'Input: "OQC meeting in the morning, NPI 407056 design in the afternoon"\n→ AI: ✓ OQC / MEET / 4h   ✓ NPI 407056 / DESIGN / 4h',
            L, Emu(4250000), LEFT_W, Emu(650000),
            bg=RGBColor(0xE8, 0xF5, 0xE9), color=RGBColor(0x1B, 0x5E, 0x20), size=10)

    # Cost classification
    textbox(slide, "Feature 2: Automated Cost Classification", L, Emu(5050000), Emu(5000000), Emu(250000),
            size=12, bold=True, color=NAVY_DARK, font=FONT_BOLD)
    textbox(slide, "System classifies into 4 cost buckets: Direct Product (60-70%) · Direct Project (10-15%) · Indirect (15-20%) · Overhead (<10%)\n95%+ accuracy · Rule engine configurable · Full audit trail",
            L, Emu(5350000), LEFT_W, Emu(600000), size=10, color=TEXT_LIGHT)

    # Right: Screenshot
    RIGHT_X = Emu(6200000)
    RIGHT_W = Emu(5600000)
    subtitle(slide, "Worklog Screen — Live System", RIGHT_X, Emu(1650000), width=RIGHT_W)
    wl = os.path.join(SCREENSHOTS_DIR, "worklogs.png")
    if os.path.exists(wl):
        slide.shapes.add_picture(wl, RIGHT_X, Emu(2000000), RIGHT_W, Emu(3400000))
    callout(slide,
            "Weekly worklog view: engineers log daily activities. Each entry captures project, activity type, hours. Data flows instantly into dashboards and cost reports.",
            RIGHT_X, Emu(5500000), RIGHT_W, Emu(650000), size=10)

    footer_line(slide)
    page_num(slide, 3)


def build_slide_4(prs, blank):
    """Management & Analytics"""
    slide = prs.slides.add_slide(blank)
    top_bar(slide)
    page_title(slide, "Core Features — Management & Analytics")

    L = Emu(500000)
    MID = Emu(6200000)
    W = Emu(5500000)

    # Left: Dashboard
    subtitle(slide, "Feature 3: Personal Dashboard", L, Emu(850000))
    db = os.path.join(SCREENSHOTS_DIR, "dashboard.png")
    if os.path.exists(db):
        slide.shapes.add_picture(db, L, Emu(1200000), W, Emu(2600000))
    callout(slide, "Login → instant overview: FTE allocation, plan vs. actual charts, overload alerts (FTE > 1.0), upcoming milestones.",
            L, Emu(3900000), W, Emu(500000), size=10)

    # Right: Resource Matrix
    subtitle(slide, "Feature 4: Resource Matrix (Org-wide)", MID, Emu(850000), width=W)
    rm = os.path.join(SCREENSHOTS_DIR, "resource-matrix.png")
    if os.path.exists(rm):
        slide.shapes.add_picture(rm, MID, Emu(1200000), W, Emu(2600000))
    callout(slide, "Complete picture: Y-axis (team members) × X-axis (months). FTE with color coding. TBD positions identify hiring gaps.",
            MID, Emu(3900000), W, Emu(500000), size=10)

    # Bottom: Forecasting
    subtitle(slide, "Feature 5: Resource Forecasting — 12-Month Forward", L, Emu(4550000), width=Emu(6500000))
    t = add_table(slide, 6, 5, L, Emu(4900000), Emu(6500000), Emu(1600000))
    t.columns[0].width = Emu(1600000)
    for c in range(1, 5):
        t.columns[c].width = Emu(1225000)
    header_row(t, ["Name", "Apr", "May", "Jun", "Status"])
    fc = [
        (["C.S. Kim", "0.8", "0.8", "0.5", "Normal"], [NAVY_DARK, TEXT, TEXT, TEXT, SUCCESS]),
        (["Y.H. Lee", "1.0", "1.0", "1.0", "Caution"], [NAVY_DARK, TEXT, TEXT, TEXT, WARNING]),
        (["J.H. Park", "0.6", "0.9", "1.2", "Alert"], [NAVY_DARK, TEXT, TEXT, DANGER, DANGER]),
        (["TBD-SW Eng", "0.5", "0.8", "1.0", "Unassigned"], [GOLD_DARK, TEXT, TEXT, TEXT, GOLD_DARK]),
        (["Gap", "-0.1", "-0.5", "-0.7", "Understaffed"], [NAVY_DARK, TEXT, DANGER, DANGER, DANGER]),
    ]
    for i, (row, colors) in enumerate(fc):
        bolds = [True] + [False] * 3 + [True]
        data_row(t, i + 1, row, size=10, colors=colors, bolds=bolds)

    # Gap callout
    callout(slide, "Jun Gap: -0.7 FTE — If hiring doesn't start now, project timelines will slip in 3 months.",
            Emu(7200000), Emu(4900000), Emu(4600000), Emu(500000),
            bg=RGBColor(0xFD, 0xED, 0xED), color=DANGER, size=11, bold=True, border_color=DANGER)

    # Feature cards
    fc_items = [("Monthly FTE", "Per project × role × month", PRIMARY),
                ("TBD Positions", "Pre-plan unhired roles", GOLD_DARK),
                ("Dynamic Capacity", "Holidays & leave adjusted", SAGE)]
    for i, (title, desc, accent) in enumerate(fc_items):
        x = Emu(7200000 + i * 1550000)
        box = add_rounded_rect(slide, x, Emu(5550000), Emu(1400000), Emu(800000), fill_color=GRAY_BG)
        add_shape(slide, x, Emu(5550000), Emu(1400000), Emu(27000), fill_color=accent)
        textbox(slide, title, x + Emu(80000), Emu(5630000), Emu(1250000), Emu(250000),
                size=11, bold=True, color=accent, font=FONT_BOLD, align=PP_ALIGN.CENTER)
        textbox(slide, desc, x + Emu(80000), Emu(5900000), Emu(1250000), Emu(300000),
                size=9, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    footer_line(slide)
    page_num(slide, 4)


def build_slide_5(prs, blank):
    """PM Core Value"""
    slide = prs.slides.add_slide(blank)
    top_bar(slide)
    page_title(slide, "Core Value for PMs — Data-Backed Resource Justification")

    L = Emu(500000)
    MID = Emu(6200000)
    W = Emu(5500000)

    # Left: Before/After
    subtitle(slide, "Before vs. After", L, Emu(850000))
    t = add_table(slide, 6, 3, L, Emu(1200000), W, Emu(2400000))
    t.columns[0].width = Emu(1100000)
    t.columns[1].width = Emu(2200000)
    t.columns[2].width = Emu(2200000)
    cell(t, 0, 0, "", size=10, bold=True, color=WHITE, bg=PRIMARY)
    cell(t, 0, 1, "Before", size=10, bold=True, color=WHITE, bg=RGBColor(0x7A, 0x2E, 0x2E), align=PP_ALIGN.CENTER)
    cell(t, 0, 2, "With Operation Board", size=10, bold=True, color=WHITE, bg=RGBColor(0x1A, 0x6B, 0x40), align=PP_ALIGN.CENTER)
    cmp = [
        ["Resource Req", '"We need more people."', '"Avg FTE 1.15, 12h OT/mo"'],
        ["Manager", '"Everyone says they\'re busy."', '"Data is solid. Start hiring."'],
        ["Data Prep", "Days of Excel cleanup", "Screen share — instant"],
        ["Budget", '"Cost breakdown?"', '"75% Direct Product"'],
        ["Quarterly", "Scramble to build report", "Auto-generated, current"],
    ]
    for i, r in enumerate(cmp):
        bg = GRAY_BG if i % 2 == 0 else WHITE
        cell(t, i + 1, 0, r[0], size=10, bold=True, color=NAVY_DARK, bg=bg)
        cell(t, i + 1, 1, r[1], size=10, color=DANGER, bg=bg)
        cell(t, i + 1, 2, r[2], size=10, color=SUCCESS, bg=bg)

    # Scenario boxes
    subtitle(slide, "The Scenario That Changes Everything", L, Emu(3800000))
    for i, (title, lines, color, bg_c) in enumerate([
        ("Without Data", 'PM: "We need 2 engineers."\nFM: "Based on what?"\nPM: "Team is stretched..."\nFM: "Next quarter."\n→ Delayed 3+ months',
         DANGER, RGBColor(0xFD, 0xED, 0xED)),
        ("With Data", 'PM: "FTE is 1.15. Jun gap: -0.7"\nFM: "Show me dashboard."\nPM: *shares screen*\nFM: "Approved. Start hiring."\n→ Approved in one meeting',
         SUCCESS, RGBColor(0xE8, 0xF5, 0xE9)),
    ]):
        x = L + Emu(i * 2850000)
        box = callout(slide, "", x, Emu(4150000), Emu(2700000), Emu(2000000), bg=bg_c)
        textbox(slide, title, x + Emu(130000), Emu(4200000), Emu(2400000), Emu(250000),
                size=12, bold=True, color=color, font=FONT_BOLD)
        textbox(slide, lines, x + Emu(130000), Emu(4450000), Emu(2400000), Emu(1600000),
                size=10, color=TEXT_LIGHT)

    # Right: Data table
    subtitle(slide, "Data PMs Can Leverage", MID, Emu(850000), width=W)
    t2 = add_table(slide, 7, 2, MID, Emu(1200000), W, Emu(2800000))
    t2.columns[0].width = Emu(2400000)
    t2.columns[1].width = Emu(3100000)
    header_row(t2, ["Data Available", "Use Case"])
    pm_data = [
        ["Planned vs Actual FTE (6mo)", "Prove overload → justify headcount"],
        ["FTE > 1.0 Heatmap", "Identify overtime concentration"],
        ["Unassigned TBD Positions", "Clarify when/what role to hire"],
        ["Cost Bucket Ratio", "Prove R&D contribution"],
        ["Capacity Gap Forecast", "Predict shortages 3-12mo ahead"],
        ["Cross-project FTE", "Identify rebalancing opportunities"],
    ]
    for i, r in enumerate(pm_data):
        data_row(t2, i + 1, r, size=10, bolds=[True, False], colors=[NAVY_DARK, TEXT])

    # Key message
    callout(slide, "Stop guessing. Start proving.\nResource requests without data get rejected.\nResource requests with data are hard to deny.",
            MID, Emu(4150000), W, Emu(900000),
            bg=RGBColor(0xE0, 0xF0, 0xF5), color=NAVY_DARK, size=13, bold=True, border_color=PRIMARY)

    # Metric cards
    metrics = [("80%", "Faster report prep"), ("6mo", "Historical trend data"), ("95%+", "Cost accuracy")]
    for i, (val, label) in enumerate(metrics):
        x = MID + Emu(i * 1900000)
        box = add_rounded_rect(slide, x, Emu(5200000), Emu(1700000), Emu(800000), fill_color=GRAY_BG)
        textbox(slide, val, x, Emu(5280000), Emu(1700000), Emu(400000),
                size=20, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER, font=FONT_BOLD)
        textbox(slide, label, x, Emu(5700000), Emu(1700000), Emu(250000),
                size=10, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    footer_line(slide)
    page_num(slide, 5)


def build_slide_6(prs, blank):
    """Organization Expansion"""
    slide = prs.slides.add_slide(blank)
    top_bar(slide)
    page_title(slide, "Organization-wide Expansion & Adoption")

    L = Emu(500000)
    MID = Emu(6200000)
    W = Emu(5500000)

    # Left: Network Effect
    subtitle(slide, "Why the Entire Org Needs This", L, Emu(850000))
    t = add_table(slide, 6, 3, L, Emu(1200000), W, Emu(2200000))
    t.columns[0].width = Emu(1200000)
    t.columns[1].width = Emu(2150000)
    t.columns[2].width = Emu(2150000)
    header_row(t, ["", "Current (4 Div)", "Full Organization"])
    ne = [
        ["Visibility", "Our team only", "Cross-functional"],
        ["Resources", "Partial teams", "All projects, all staff"],
        ["FTE Data", "4 divisions only", "Org-wide capacity"],
        ["Forecasting", "Our team only", "Org-wide supply/demand"],
        ["Decisions", "Team level", "Program / Org level"],
    ]
    for i, r in enumerate(ne):
        bg = GRAY_BG if i % 2 == 0 else WHITE
        cell(t, i + 1, 0, r[0], size=10, bold=True, color=NAVY_DARK, bg=bg)
        cell(t, i + 1, 1, r[1], size=10, color=GRAY, bg=bg)
        cell(t, i + 1, 2, r[2], size=10, bold=True, color=SUCCESS, bg=bg)

    callout(slide,
            "NPI projects involve Control, Mechanical, and Electrical teams.\nOne team = 1/3 of the picture. All teams = complete project resource view.",
            L, Emu(3500000), W, Emu(600000),
            bg=RGBColor(0xE8, 0xF5, 0xE9), color=RGBColor(0x1B, 0x5E, 0x20), size=11, bold=True)

    # Coverage visual
    subtitle(slide, "Data Coverage Impact", L, Emu(4250000))
    textbox(slide, "Today: 4/8 Divisions  ▶▶▶  Goal: 8/8 Divisions",
            L, Emu(4600000), W, Emu(350000), size=12, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    textbox(slide, "Partial visibility → Complete program-level optimization",
            L, Emu(4950000), W, Emu(300000), size=11, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    # Right: Adoption Burden
    subtitle(slide, "Adoption Burden: Virtually None", MID, Emu(850000), width=W)
    t2 = add_table(slide, 6, 2, MID, Emu(1200000), W, Emu(2000000))
    t2.columns[0].width = Emu(1500000)
    t2.columns[1].width = Emu(4000000)
    header_row(t2, ["Item", "Detail"])
    ab = [
        ["Cost", "None — self-hosted, no license"],
        ["Installation", "None — web browser only"],
        ["Training", "Under 30 min — worklog + dashboard"],
        ["Daily Effort", "30 seconds/day"],
        ["Data Migration", "Auto-migration from SharePoint"],
    ]
    for i, r in enumerate(ab):
        data_row(t2, i + 1, r, size=10, bolds=[True, False], colors=[NAVY_DARK, TEXT])

    # Timeline
    subtitle(slide, "Rollout Timeline", MID, Emu(3350000), width=W)
    timeline = [("Week 1", "Account creation\n30-min onboarding", PRIMARY),
                ("Week 2-3", "Daily worklogs\n30 sec/day", NAVY),
                ("Week 4+", "Dashboard insights\nResource planning", SAGE)]
    for i, (week, desc, accent) in enumerate(timeline):
        x = MID + Emu(i * 1900000)
        card = add_rounded_rect(slide, x, Emu(3700000), Emu(1700000), Emu(1000000), fill_color=accent)
        textbox(slide, week, x, Emu(3760000), Emu(1700000), Emu(300000),
                size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER, font=FONT_BOLD)
        textbox(slide, desc, x + Emu(100000), Emu(4050000), Emu(1500000), Emu(550000),
                size=10, color=WHITE, align=PP_ALIGN.CENTER)

    # FAQ
    subtitle(slide, "FAQ", MID, Emu(4850000), width=W)
    faqs = [
        ('Q: "Will this add more work?"', "A: 30 sec/day — less than one email."),
        ('Q: "What if data isn\'t perfect?"', "A: 70% coverage already beats spreadsheets."),
        ('Q: "Can we go back to SharePoint?"', "A: Yes. But none of the 4 pilots have."),
    ]
    for i, (q, a) in enumerate(faqs):
        y = Emu(5150000 + i * 420000)
        box = add_rounded_rect(slide, MID, y, W, Emu(360000), fill_color=GRAY_BG)
        textbox(slide, q, MID + Emu(120000), y + Emu(30000), Emu(5200000), Emu(170000),
                size=10, bold=True, color=NAVY_DARK, font=FONT_BOLD)
        textbox(slide, a, MID + Emu(120000), y + Emu(190000), Emu(5200000), Emu(160000),
                size=10, color=TEXT_LIGHT)

    footer_line(slide)
    page_num(slide, 6)


def build_slide_7(prs, blank):
    """In-House Capability"""
    slide = prs.slides.add_slide(blank)
    top_bar(slide)
    page_title(slide, "In-House Capability & Operational Results")

    L = Emu(500000)
    MID = Emu(6200000)
    W = Emu(5500000)

    # Left: SaaS vs In-house
    subtitle(slide, "External SaaS vs. In-House Built", L, Emu(850000))
    t = add_table(slide, 7, 3, L, Emu(1200000), W, Emu(2600000))
    t.columns[0].width = Emu(1100000)
    t.columns[1].width = Emu(2200000)
    t.columns[2].width = Emu(2200000)
    header_row(t, ["", "External SaaS", "Operation Board"])
    saas = [
        ["Cost", "License/user — millions KRW/yr", "Self-hosted, no cost"],
        ["Custom", "Limited, vendor-dependent", "100% tailored, same-day"],
        ["Data", "Stored on vendor cloud", "Our servers, full control"],
        ["Speed", "Weeks to months per request", "Immediate"],
        ["Process", "Force-fit generic features", "Edwards PCP native"],
        ["Korea", "No Korean holiday support", "Korean holidays built-in"],
    ]
    for i, r in enumerate(saas):
        bg = GRAY_BG if i % 2 == 0 else WHITE
        cell(t, i + 1, 0, r[0], size=10, bold=True, color=NAVY_DARK, bg=bg)
        cell(t, i + 1, 1, r[1], size=10, color=TEXT_LIGHT, bg=bg)
        cell(t, i + 1, 2, r[2], size=10, bold=True, color=SUCCESS, bg=bg)

    callout(slide,
            "Designed and built in-house by our engineering organization. Tailored to our PCP process, org structure, and project types.",
            L, Emu(3950000), W, Emu(500000),
            bg=RGBColor(0xE0, 0xF0, 0xF5), color=PRIMARY, size=11, bold=True, border_color=PRIMARY)

    # Right: Results
    subtitle(slide, "Operational Results — Before & After", MID, Emu(850000), width=W)
    t2 = add_table(slide, 7, 4, MID, Emu(1200000), W, Emu(2600000))
    t2.columns[0].width = Emu(1300000)
    t2.columns[1].width = Emu(1300000)
    t2.columns[2].width = Emu(1400000)
    t2.columns[3].width = Emu(1500000)
    header_row(t2, ["Metric", "Before", "After", "Improvement"])
    res = [
        ["Worklog", "2-3 min", "Under 30 sec", "80% faster"],
        ["Resource", "Half day", "Real-time", "Instant"],
        ["Report", "1-2 days", "Auto-gen", "Automated"],
        ["Forecast", "Impossible", "12-month", "New"],
        ["Cross-proj", "Individual", "Single screen", "Real-time"],
        ["Cost Class", "Manual", "Auto 95%+", "Automated"],
    ]
    for i, r in enumerate(res):
        bg = GRAY_BG if i % 2 == 0 else WHITE
        cell(t2, i + 1, 0, r[0], size=10, bold=True, color=NAVY_DARK, bg=bg)
        cell(t2, i + 1, 1, r[1], size=10, color=DANGER, bg=bg)
        cell(t2, i + 1, 2, r[2], size=10, color=SUCCESS, bg=bg)
        cell(t2, i + 1, 3, r[3], size=10, bold=True, color=PRIMARY, bg=bg, align=PP_ALIGN.CENTER)

    # Tech stack
    subtitle(slide, "Technology Stack", MID, Emu(3950000), width=W)
    stacks = [
        ("Frontend", "React + TypeScript\nTanStack Query\nTailwind CSS", PRIMARY),
        ("Backend", "FastAPI (Python)\nSQLAlchemy ORM\nJWT Auth", NAVY),
        ("Infra", "PostgreSQL\nDocker\nSelf-hosted", SAGE),
        ("AI", "NLP Parser\nRule Classifier\nDashboards", GOLD_DARK),
    ]
    for i, (title, desc, accent) in enumerate(stacks):
        x = MID + Emu(i * 1400000)
        box = add_rounded_rect(slide, x, Emu(4300000), Emu(1300000), Emu(1200000), fill_color=GRAY_BG)
        add_shape(slide, x, Emu(4300000), Emu(1300000), Emu(27000), fill_color=accent)
        textbox(slide, title, x + Emu(80000), Emu(4380000), Emu(1150000), Emu(250000),
                size=11, bold=True, color=accent, font=FONT_BOLD)
        textbox(slide, desc, x + Emu(80000), Emu(4650000), Emu(1150000), Emu(700000),
                size=9, color=TEXT_LIGHT)

    callout(slide, "Modern, production-grade architecture. Same tech as leading SaaS companies. Fully containerized. No vendor lock-in.",
            L, Emu(5600000), Emu(11200000), Emu(400000), size=10, color=TEXT_LIGHT)

    footer_line(slide)
    page_num(slide, 7)


def build_slide_8(prs, blank):
    """Roadmap & Proposal"""
    slide = prs.slides.add_slide(blank)
    top_bar(slide)
    page_title(slide, "Roadmap & Proposal")

    L = Emu(500000)
    MID = Emu(6200000)
    W = Emu(5500000)

    # Left: Roadmap
    subtitle(slide, "Development Roadmap", L, Emu(850000))
    phases = [
        ("Phase 1", "Completed", "Worklog, Resource matrix, Dashboard, PCP milestones, Cost classification, AI input", PRIMARY, SUCCESS),
        ("Phase 2", "3-6 months", "Team dashboard, FTE>1.0 alerts, Mobile UI, Power BI integration", NAVY, GOLD_DARK),
        ("Phase 3", "6-12 months", "What-if simulation, Pattern analytics, SAP integration, Org optimization", SAGE, GOLD_DARK),
        ("Phase 4", "12+ months", "AI resource recommendations, Risk detection, Global multi-site expansion", GRAY, TEXT_LIGHT),
    ]
    for i, (phase, timing, desc, accent, timing_color) in enumerate(phases):
        y = Emu(1200000 + i * 750000)
        box = add_rounded_rect(slide, L, y, W, Emu(650000), fill_color=GRAY_BG)
        add_shape(slide, L, y, Emu(35000), Emu(650000), fill_color=accent)
        textbox(slide, phase, L + Emu(100000), y + Emu(50000), Emu(900000), Emu(250000),
                size=12, bold=True, color=accent, font=FONT_BOLD)
        textbox(slide, timing, L + Emu(100000), y + Emu(300000), Emu(900000), Emu(200000),
                size=9, bold=True, color=timing_color)
        textbox(slide, desc, L + Emu(1100000), y + Emu(100000), Emu(4200000), Emu(450000),
                size=10, color=TEXT)

    callout(slide,
            "Key advantage: Every phase driven by direct user feedback — not a vendor's roadmap. Features prioritized by what our PMs and FMs need.",
            L, Emu(4300000), W, Emu(500000),
            bg=RGBColor(0xFE, 0xF3, 0xE2), color=GOLD_DARK, size=10, bold=True, border_color=GOLD)

    # Right: Next Steps
    subtitle(slide, "Next Steps by Role", MID, Emu(850000), width=W)
    t = add_table(slide, 4, 3, MID, Emu(1200000), W, Emu(1300000))
    t.columns[0].width = Emu(1000000)
    t.columns[1].width = Emu(3300000)
    t.columns[2].width = Emu(1200000)
    header_row(t, ["Role", "Ask", "Effort"])
    steps = [
        ["Divisions", "Join 2-week pilot → start logging → experience results", "30 sec/day"],
        ["PMs", "Register Excel resource plan (one-time) → real-time tracking", "One-time"],
        ["FMs", "Review FTE dashboard weekly, manage TBD positions", "Weekly"],
    ]
    for i, r in enumerate(steps):
        data_row(t, i + 1, r, size=10, bolds=[True, False, True], colors=[NAVY_DARK, TEXT, SUCCESS])

    # Support
    subtitle(slide, "Support We Provide", MID, Emu(2650000), width=W)
    support = [
        "30-minute onboarding session per team",
        "Auto-migration from existing SharePoint data",
        "Continuous improvement — feature requests deployed within days",
        "Dedicated support channel — direct access to dev team",
    ]
    sup_box = slide.shapes.add_textbox(MID, Emu(2950000), W, Emu(1200000))
    tf = sup_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(support):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT
        p.font.name = FONT
        p.space_before = Pt(4)

    # Closing message
    closing = add_rounded_rect(slide, MID, Emu(4200000), W, Emu(800000), fill_color=PRIMARY)
    box = slide.shapes.add_textbox(MID + Emu(200000), Emu(4280000), Emu(5100000), Emu(650000))
    tf = set_text(box, "This is not mandatory. Just try it for 2 weeks.", size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_BOLD)
    add_para(tf, "No risk. Go back if it doesn't work.", size=11, color=RGBColor(0xCC, 0xDD, 0xDD), align=PP_ALIGN.CENTER)
    add_para(tf, "However, none of the 4 pilot divisions have gone back yet.", size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Footer bar
    add_shape(slide, Emu(0), Emu(6500000), SLIDE_W, Emu(358000), fill_color=PRIMARY)
    textbox(slide, "Access: System URL   |   Demo: guest@edwards.com   |   Contact: EUV Program IS",
            Emu(0), Emu(6540000), SLIDE_W, Emu(250000),
            size=10, color=RGBColor(0xCC, 0xDD, 0xDD), align=PP_ALIGN.CENTER)


def main():
    # Use AC template
    if os.path.exists(TEMPLATE_PATH):
        prs = Presentation(TEMPLATE_PATH)
        print(f"Using template: {TEMPLATE_PATH}")
    else:
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H
        print("Template not found, using blank presentation")

    # Find Blank layout
    blank = None
    for layout in prs.slide_layouts:
        if layout.name == "Blank":
            blank = layout
            break
    if blank is None:
        blank = prs.slide_layouts[0]
        print(f"Using layout: '{blank.name}' (Blank not found)")
    else:
        print(f"Using layout: '{blank.name}'")

    # Remove template sample slides
    from pptx.oxml.ns import qn
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn('r:id'))
        if rId:
            prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)
    print(f"Cleared template slides. Building editable slides...")

    print("  Slide 1: Cover")
    build_slide_1(prs, blank)
    print("  Slide 2: Background & Design Principles")
    build_slide_2(prs, blank)
    print("  Slide 3: Input Features")
    build_slide_3(prs, blank)
    print("  Slide 4: Management & Analytics")
    build_slide_4(prs, blank)
    print("  Slide 5: PM Core Value")
    build_slide_5(prs, blank)
    print("  Slide 6: Organization Expansion")
    build_slide_6(prs, blank)
    print("  Slide 7: In-House Capability")
    build_slide_7(prs, blank)
    print("  Slide 8: Roadmap & Proposal")
    build_slide_8(prs, blank)

    prs.save(OUTPUT_PATH)
    print(f"\nPPT saved: {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")
    print(f"File size: {os.path.getsize(OUTPUT_PATH) / 1024 / 1024:.1f} MB")
    print("All elements are EDITABLE in PowerPoint!")


if __name__ == "__main__":
    main()
